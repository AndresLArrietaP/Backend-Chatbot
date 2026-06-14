"""
Genera CONFIA_Presentacion_Gerencia.pptx
Presentación ejecutiva para vender el proyecto KomfIA / CONFIA a gerencia KMMP.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── Colores corporativos ────────────────────────────────────────────────────
KOM_BLUE  = RGBColor(0x00, 0x30, 0x87)
KOM_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xEC, 0xF2, 0xFB)
GRAY      = RGBColor(0x55, 0x55, 0x65)
RED_ACC   = RGBColor(0xCC, 0x00, 0x00)
GREEN_ACC = RGBColor(0x00, 0x7A, 0x33)
ORANGE    = RGBColor(0xE8, 0x8A, 0x1A)
TEAL      = RGBColor(0x00, 0x6E, 0x8A)

# ─── Presentación ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, l, t, w, h, fill=LIGHT_BG, line=None, line_w=Pt(1)):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp


def txt(slide, l, t, w, h, text, size=12, bold=False, color=KOM_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_txt(slide, l, t, w, h, lines, base_size=12, line_gap=3):
    """lines = list of (text, size, bold, color, align)"""
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = Pt(line_gap)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def bullet_block(slide, l, t, w, h, title, bullets,
                 title_size=13, bullet_size=11, bg=LIGHT_BG,
                 title_color=KOM_BLUE, bullet_color=KOM_DARK,
                 border_color=KOM_BLUE):
    box(slide, l, t, w, h, fill=bg, line=border_color, line_w=Pt(1.2))
    txt(slide, l+0.12, t+0.1, w-0.24, 0.42,
        title, title_size, bold=True, color=title_color)
    tb = slide.shapes.add_textbox(
        Inches(l+0.12), Inches(t+0.55), Inches(w-0.24), Inches(h-0.68)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = "• " + b
        r.font.size = Pt(bullet_size)
        r.font.color.rgb = bullet_color


def header_bar(slide, title, subtitle=None):
    """Barra azul superior con título blanco."""
    box(slide, 0, 0, 13.33, 1.1, fill=KOM_BLUE)
    txt(slide, 0.4, 0.1, 12.5, 0.7, title,
        26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, 0.4, 0.72, 12.5, 0.35, subtitle,
            11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)


def footer(slide, num=None):
    box(slide, 0, 7.18, 13.33, 0.32, fill=KOM_BLUE)
    txt(slide, 0.25, 7.19, 3.5, 0.28,
        "KOMATSU  |  MITSUI", 9, bold=True, color=WHITE)
    txt(slide, 4.5, 7.19, 4.33, 0.28,
        "Confidential", 9, bold=True, color=RGBColor(0xFF, 0xDD, 0xAA),
        align=PP_ALIGN.CENTER)
    right = "Komatsu - Mitsui Maquinarias Perú S.A."
    if num:
        right += f"    {num}"
    txt(slide, 9.0, 7.19, 4.1, 0.28, right, 9, color=WHITE, align=PP_ALIGN.RIGHT)


def logo_text(slide, l=0.35, t=0.15):
    txt(slide, l, t, 3.0, 0.55, "KOMATSU", 22, bold=True, color=KOM_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s1 = blank(prs)

# Franja azul izquierda
box(s1, 0, 0, 0.15, 7.5, fill=KOM_BLUE)

# Logo
logo_text(s1, 0.45, 0.25)

# Línea superior decorativa
box(s1, 0.45, 0.9, 11.0, 0.04, fill=KOM_BLUE)

# Título principal
txt(s1, 0.45, 1.3, 10.0, 1.8,
    "CONFIA", 82, bold=True, color=KOM_BLUE, align=PP_ALIGN.LEFT)

# Subtítulo
txt(s1, 0.45, 3.15, 10.5, 0.7,
    "Plataforma de Inteligencia Artificial para Gestión de Confiabilidad",
    19, color=GRAY, align=PP_ALIGN.LEFT)

# Separador
box(s1, 0.45, 3.95, 8.5, 0.05, fill=KOM_BLUE)

# Autor
txt(s1, 0.45, 4.15, 9.0, 0.5,
    "Andrés Arrieta Palacios", 18, bold=True, color=KOM_DARK)
txt(s1, 0.45, 4.68, 9.0, 0.4,
    "Practicante Pre Profesional  —  Área OIS  |  Komatsu Mitsui Maquinarias Perú S.A.",
    12, color=GRAY)

# Badge
box(s1, 0.45, 5.3, 3.2, 0.7, fill=KOM_BLUE)
txt(s1, 0.55, 5.38, 3.0, 0.52,
    "PRESENTACIÓN EJECUTIVA", 11, bold=True, color=WHITE)

footer(s1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s2 = blank(prs)
header_bar(s2, "El Problema: Monitoreo Reactivo y Manual",
           "El análisis de aceite es el corazón de la confiabilidad — pero hoy depende de un proceso 100% manual.")

problems = [
    ("⏱  Proceso Manual",
     ["Revisión muestra por muestra por el especialista",
      "Horas de trabajo por ciclo de análisis",
      "Escalabilidad limitada: depende de 1 persona",
      "Vacaciones o ausencias = proceso paralizado"]),
    ("📅  Información No Oportuna",
     ["Reportes entregados de forma diaria, semanal o mensual",
      "Cuando llega el informe, la falla ya pudo ocurrir",
      "Sin alertas inmediatas ni priorización automática",
      "El monitoreo de condición pierde su valor preventivo"]),
    ("🚨  Consecuencias Directas",
     ["Sin acción oportuna → falla catastrófica del componente",
      "KMT: reparación estimada en S/. 1.8M por evento",
      "Camión paralizado: 2+ días fuera de operación",
      "Impacto en disponibilidad, costos y seguridad"]),
]
for i, (title, bullets) in enumerate(problems):
    bullet_block(s2, 0.38 + i * 4.25, 1.25, 4.05, 5.55,
                 title, bullets, title_size=14, bullet_size=12)

footer(s2, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — IMPACTO EN NÚMEROS
# ══════════════════════════════════════════════════════════════════════════════
s3 = blank(prs)
header_bar(s3, "El Costo de No Actuar a Tiempo",
           "Números reales del impacto cuando el monitoreo no es oportuno.")

metrics = [
    ("S/ 1.8 M",    "Reparación KMT por componente\n(Motor de Tracción, Rueda, Hidráulico)", KOM_BLUE,  WHITE),
    ("S/ 3.6 M+",   "Impacto total por camión paralizado\n(componente + lucro cesante estimado)", KOM_BLUE,  WHITE),
    ("↓ Disponibilidad", "Cada hora parada impacta directamente\nel KPI de disponibilidad de flota",      LIGHT_BG, KOM_BLUE),
    ("↑ Riesgo",    "Operar con un componente en estado\ndesconocido es un riesgo de seguridad",          LIGHT_BG, KOM_BLUE),
]
for i, (value, desc, bg, fc) in enumerate(metrics):
    col, row = i % 2, i // 2
    lx = 0.4 + col * 6.35
    ty = 1.45 + row * 2.65
    box(s3, lx, ty, 6.05, 2.4, fill=bg,
        line=KOM_BLUE, line_w=Pt(2))
    txt(s3, lx+0.18, ty+0.15, 5.7, 0.9,
        value, 34, bold=True, color=fc)
    txt(s3, lx+0.18, ty+1.05, 5.7, 1.2,
        desc, 12.5, color=fc if bg == KOM_BLUE else GRAY)

footer(s3, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LA SOLUCIÓN: KOMFIA / CONFIA
# ══════════════════════════════════════════════════════════════════════════════
s4 = blank(prs)
header_bar(s4, "La Solución: KomfIA",
           "Un asistente de IA que convierte preguntas en lenguaje natural en análisis reales, en segundos.")

# Descripción central
box(s4, 0.38, 1.25, 12.57, 1.05, fill=LIGHT_BG)
txt(s4, 0.55, 1.35, 12.2, 0.85,
    'El usuario escribe en Teams: "¿Qué motores de tracción de Antapaccay están observados?" — '
    'y KomfIA consulta la base de datos, evalúa los límites reales LP/LC y responde al instante.',
    13.5, color=KOM_DARK, italic=True)

feats = [
    ("💬\nLenguaje Natural",
     "Sin SQL, sin reportes.\nEl usuario habla, KomfIA entiende."),
    ("⚡\nRespuesta Inmediata",
     "Resultados en segundos.\nNo en horas ni días."),
    ("🔗\nDatos Reales",
     "Consulta directa a la BD\noperativa de KMMP. Sin intermediarios."),
    ("📊\nAnálisis Enriquecido",
     "Tendencias, triage, rankings\ny recomendaciones técnicas."),
]
for i, (title, desc) in enumerate(feats):
    lx = 0.38 + i * 3.15
    box(s4, lx, 2.5, 3.0, 4.2,
        fill=KOM_BLUE if i == 0 else LIGHT_BG,
        line=KOM_BLUE, line_w=Pt(1.5))
    tc = WHITE if i == 0 else KOM_BLUE
    dc = WHITE if i == 0 else GRAY
    txt(s4, lx+0.12, 2.65, 2.76, 1.1, title, 13, bold=True,
        color=tc, align=PP_ALIGN.CENTER)
    txt(s4, lx+0.12, 3.85, 2.76, 2.6, desc, 12, color=dc,
        align=PP_ALIGN.CENTER)

footer(s4, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FLUJO DE FUNCIONAMIENTO
# ══════════════════════════════════════════════════════════════════════════════
s5 = blank(prs)
header_bar(s5, "Flujo de Funcionamiento",
           "Del lenguaje natural a la respuesta enriquecida en 4 pasos.")

# ── Nodos del flujo ───────────────────────────────────────────────────────────
NW, NH = 2.3, 1.15

nodes = [
    ("USUARIO",       "Teams / Copilot Studio", 0.35,  3.15, KOM_BLUE),
    ("BACKEND",       "Python  (FastAPI)",       3.3,   2.2,  KOM_BLUE),
    ("LLM",           "Gemini / OpenAI",         5.95,  1.25, TEAL),
    ("BASE DE DATOS", "Azure SQL Server",        5.95,  4.4,  TEAL),
    ("RESPUESTA",     "JSON enriquecido\nal agente", 9.25, 3.15, GREEN_ACC),
]
for label, sub, lx, ty, fc in nodes:
    box(s5, lx, ty, NW, NH, fill=fc)
    txt(s5, lx+0.08, ty+0.1,  NW-0.16, 0.48,
        label, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s5, lx+0.08, ty+0.58, NW-0.16, 0.48,
        sub, 9.5, color=WHITE, align=PP_ALIGN.CENTER)

# ── Conectores (flechas simplificadas como barras + etiquetas) ─────────────────
# ① USUARIO → BACKEND  (horizontal)
box(s5, 2.65, 3.67, 0.65, 0.06, fill=KOM_BLUE)
txt(s5, 2.35, 3.41, 1.2, 0.28, "① Consulta NL", 8.5, bold=True, color=KOM_BLUE, align=PP_ALIGN.CENTER)

# ② BACKEND → LLM  (diagonal up-right: usar barra vertical + horizontal)
box(s5, 4.6,  2.6,  1.35, 0.06, fill=TEAL)   # horizontal
box(s5, 5.9,  1.35, 0.06, 1.3,  fill=TEAL)   # vertical
txt(s5, 4.35, 2.38, 2.0, 0.25, "② Traduce a SQL (LLM)", 8.5, color=TEAL)

# ③ LLM → BACKEND  (SQL retornado)
box(s5, 4.6,  2.75, 1.35, 0.06, fill=KOM_BLUE)   # horizontal (ligeramente abajo)
txt(s5, 4.35, 2.82, 2.0, 0.25, "③ SQL generado", 8.5, color=KOM_BLUE)

# ④ BACKEND → BD  (vertical down)
box(s5, 4.43, 3.35, 0.06, 1.05, fill=TEAL)
txt(s5, 3.55, 4.05, 1.55, 0.25, "④ Ejecuta SQL", 8.5, color=TEAL)

# ⑤ BD → BACKEND  (data)
box(s5, 4.55, 4.45, 1.4, 0.06, fill=KOM_BLUE)
txt(s5, 4.35, 4.52, 1.8, 0.25, "⑤ Retorna datos", 8.5, color=KOM_BLUE)

# ⑥ BACKEND → RESPUESTA  (horizontal right)
box(s5, 7.62, 3.67, 1.63, 0.06, fill=GREEN_ACC)
txt(s5, 7.25, 3.41, 2.4, 0.28, "⑥ Respuesta JSON + análisis", 8.5, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)

# Nota de contexto
box(s5, 0.35, 6.55, 12.57, 0.52, fill=LIGHT_BG)
txt(s5, 0.5, 6.6, 12.3, 0.43,
    "El backend también guarda contexto conversacional — el agente recuerda la sesión para refinamientos "
    "como \"de esos, muéstrame el peor\" o \"quédate con los críticos\" sin repetir filtros.",
    10.5, color=KOM_BLUE, italic=True)

footer(s5, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CAPACIDADES ACTUALES
# ══════════════════════════════════════════════════════════════════════════════
s6 = blank(prs)
header_bar(s6, "¿Qué puede hacer KomfIA hoy?",
           "Capacidades operativas activas — probadas en marcha blanca con flota Antapaccay.")

caps = [
    ("🔎  Triage de Flota",
     ["¿Qué motores de tracción están observados?",
      "Evalúa el último análisis de cada equipo",
      "Compara vs. límites LP/LC reales de la BD",
      "Lista priorizada: CRÍTICOS primero"]),
    ("📈  Tendencias",
     ["¿Cómo ha evolucionado el Fe del CA3164?",
      "Últimas 8 muestras individuales (default)",
      "O promedio mensual / trimestral / semestral",
      "Con líneas de referencia LP/LC"]),
    ("📋  Historial Completo",
     ["Últimas N muestras de un equipo/componente",
      "Registro con fecha, metales, tipo de muestra",
      "Identifica cambios de aceite y post-dializados"]),
    ("🏆  Rankings y Comparativas",
     ["Top 5 motores con mayor Fe de Antapaccay",
      "¿Cuál es el hidráulico con más cobre?",
      "Comparativas entre equipos y proyectos"]),
]
for i, (title, bullets) in enumerate(caps):
    col, row = i % 2, i // 2
    lx = 0.38 + col * 6.4
    ty = 1.28 + row * 2.9
    bullet_block(s6, lx, ty, 6.1, 2.68,
                 title, bullets, title_size=14, bullet_size=12)

footer(s6, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — IMPACTO EN KPIs
# ══════════════════════════════════════════════════════════════════════════════
s7 = blank(prs)
header_bar(s7, "Impacto en Indicadores Clave (KPIs)",
           "La detección oportuna de anomalías impacta directamente en los indicadores que gerencia monitorea.")

kpis = [
    ("📈 Disponibilidad",          "↑  SUBE",
     "Menos fallas no planificadas = más horas de operación efectiva. "
     "Cada componente intervenido antes de fallar es un turno menos paralizado.",
     GREEN_ACC),
    ("🔧 MTBF",                    "↑  SUBE",
     "Mayor tiempo entre fallas al actuar preventivamente. "
     "El análisis oportuno permite planificar el cambio antes del daño catastrófico.",
     GREEN_ACC),
    ("💰 Costo de Mantenimiento",  "↓  BAJA",
     "Reparación preventiva vs. reparación catastrófica: diferencia de millones. "
     "KMT a S/. 1.8M — la detección temprana evita llegar a ese punto.",
     GREEN_ACC),
    ("🛡  Seguridad",              "↑  MEJORA",
     "Eliminar la operación de equipos con componentes en estado crítico desconocido. "
     "Un equipo observado y no atendido es un riesgo para el operador.",
     GREEN_ACC),
]
for i, (kpi, arrow, desc, color) in enumerate(kpis):
    col, row = i % 2, i // 2
    lx = 0.38 + col * 6.4
    ty = 1.32 + row * 2.75
    box(s7, lx, ty, 6.1, 2.5,
        fill=LIGHT_BG, line=color, line_w=Pt(2.5))
    txt(s7, lx+0.15, ty+0.12, 3.6, 0.6,
        kpi, 14, bold=True, color=KOM_BLUE)
    txt(s7, lx+3.9,  ty+0.08, 2.0, 0.65,
        arrow, 22, bold=True, color=color, align=PP_ALIGN.RIGHT)
    # separador
    box(s7, lx+0.15, ty+0.72, 5.8, 0.04, fill=color)
    txt(s7, lx+0.15, ty+0.82, 5.8, 1.55,
        desc, 11.5, color=GRAY)

footer(s7, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VISIÓN: MÁS ALLÁ DEL ANÁLISIS DE ACEITE
# ══════════════════════════════════════════════════════════════════════════════
s8 = blank(prs)
header_bar(s8, "Visión: CONFIA no se limita al Análisis de Aceite",
           "La plataforma es el núcleo de inteligencia operacional del área — con potencial de expansión.")

modules = [
    ("🛢  Análisis de Aceite (KomfIA)",
     "EN PRODUCCIÓN",
     "Triage, tendencias, historial, rankings.\nActivo en Antapaccay. Escalando a todos los proyectos KMMP.",
     KOM_BLUE, True),
    ("📊  Monitoreo de KPIs",
     "ROADMAP",
     "Detección automática de desviaciones en disponibilidad,\nMTBF, costos y productividad de flota.",
     ORANGE, False),
    ("🤖  PCA / Machine Learning",
     "ROADMAP",
     "Análisis de componentes principales para detección\nmultivariable de anomalías y predicción de fallas.",
     ORANGE, False),
    ("📡  Integración Invertex",
     "ROADMAP",
     "Datos de telemetría y condición operativa en\ntiempo real integrados al análisis.",
     TEAL, False),
    ("📁  Data PLM",
     "ROADMAP",
     "Historial de mantenimiento, planes y órdenes\nde trabajo como contexto del análisis.",
     TEAL, False),
    ("🚨  Alertas Automáticas",
     "ROADMAP",
     "Notificaciones por Teams / email cuando un\nequipo supera límites o muestra tendencia crítica.",
     GREEN_ACC, False),
]
for i, (title, status, desc, color, active) in enumerate(modules):
    col, row = i % 3, i // 3
    lx = 0.38 + col * 4.28
    ty = 1.28 + row * 2.7
    bg = LIGHT_BG if active else WHITE
    box(s8, lx, ty, 4.1, 2.48, fill=bg, line=color, line_w=Pt(2.0))
    txt(s8, lx+0.12, ty+0.1,  2.7, 0.42, title, 12, bold=True, color=KOM_BLUE)
    badge_bg = KOM_BLUE if active else LIGHT_BG
    badge_fc = WHITE if active else color
    box(s8, lx+2.95, ty+0.1, 1.0, 0.36, fill=badge_bg)
    txt(s8, lx+2.97, ty+0.11, 0.98, 0.32, status, 7.5, bold=True, color=badge_fc, align=PP_ALIGN.CENTER)
    txt(s8, lx+0.12, ty+0.62, 3.86, 1.75, desc, 10.5, color=GRAY)

footer(s8, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONFIA COMO MARCA OIS
# ══════════════════════════════════════════════════════════════════════════════
s9 = blank(prs)
header_bar(s9, "CONFIA: La Marca de Inteligencia del Área OIS",
           "No es solo un chatbot — es la infraestructura de automatización e inteligencia del área.")

# Nodo central
box(s9, 5.0, 1.95, 3.33, 3.3, fill=KOM_BLUE)
txt(s9, 5.0, 2.2,  3.33, 1.4, "CONFIA", 42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s9, 5.0, 3.65, 3.33, 0.5, "Sistema IA — Área OIS", 11, color=WHITE, align=PP_ALIGN.CENTER)
txt(s9, 5.0, 4.15, 3.33, 0.8, "KMMP", 11, bold=True, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Satélites
satellites = [
    ("🛢 Análisis de Aceite\n(KomfIA)  ✅",   0.3,  2.0, KOM_BLUE),
    ("📊 Monitoreo de KPIs",                    0.3,  4.5, ORANGE),
    ("🤖 PCA / Anomalías ML",                  10.5,  2.0, TEAL),
    ("📁 Invertex + PLM",                       10.5,  4.5, TEAL),
]
SAT_W, SAT_H = 2.2, 1.05
for label, lx, ty, color in satellites:
    box(s9, lx, ty, SAT_W, SAT_H, fill=LIGHT_BG, line=color, line_w=Pt(2))
    txt(s9, lx+0.1, ty+0.1, SAT_W-0.2, SAT_H-0.15,
        label, 11, bold=True, color=KOM_BLUE)

# Líneas conectoras
box(s9, 2.5,  2.48, 2.5,  0.06, fill=KOM_BLUE)
box(s9, 2.5,  4.98, 2.5,  0.06, fill=KOM_BLUE)
box(s9, 8.33, 2.48, 2.17, 0.06, fill=KOM_BLUE)
box(s9, 8.33, 4.98, 2.17, 0.06, fill=KOM_BLUE)

# Tagline
box(s9, 0.38, 6.3, 12.57, 0.65, fill=LIGHT_BG)
txt(s9, 0.55, 6.37, 12.2, 0.55,
    '"OIS tiene un sistema de inteligencia artificial llamado CONFIA que automatiza procesos, '
    'detecta anomalías y amplifica la capacidad técnica del área — '
    'hoy en análisis de aceite, mañana en KPIs, PLM e integración operacional."',
    11.5, italic=True, color=KOM_BLUE, align=PP_ALIGN.CENTER)

footer(s9, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════════════
s10 = blank(prs)
header_bar(s10, "Próximos Pasos",
           "Hoja de ruta para escalar CONFIA a toda la operación KMMP.")

phases = [
    ("CORTO PLAZO\n(Próximas 4 semanas)",
     KOM_BLUE, True,
     ["Extensión a todos los proyectos activos:\nAntamina, Cerro Verde, Toromocho",
      "Finalizar pruebas con usuarios en Antapaccay",
      "Motor de recomendaciones técnicas por metal",
      "Onboarding del equipo OIS como usuarios"]),
    ("MEDIANO PLAZO\n(1 — 3 meses)",
     ORANGE, False,
     ["Módulo de monitoreo de KPIs operacionales",
      "Alertas automáticas por Teams cuando\nun equipo supera límites",
      "Integración con datos Invertex",
      "Dashboard visual de estado de flota"]),
    ("LARGO PLAZO\n(3 — 6 meses)",
     TEAL, False,
     ["PCA para detección multivariable de fallas",
      "Integración PLM (planes y OT de mantenimiento)",
      "Extensión a otras áreas / contratos KMMP",
      "CONFIA como producto de área OIS"]),
]
for i, (phase, color, active, steps) in enumerate(phases):
    lx = 0.38 + i * 4.28
    bg = LIGHT_BG if active else WHITE
    box(s10, lx, 1.28, 4.1, 5.5, fill=bg, line=color, line_w=Pt(2.5))
    # Phase header
    box(s10, lx, 1.28, 4.1, 0.75, fill=color)
    txt(s10, lx+0.12, 1.32, 3.86, 0.65,
        phase, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s10.shapes.add_textbox(
        Inches(lx+0.15), Inches(2.15), Inches(3.8), Inches(4.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for j, step in enumerate(steps):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = "• " + step
        r.font.size = Pt(11.5)
        r.font.color.rgb = KOM_DARK

footer(s10, 10)


# ─── Guardar ──────────────────────────────────────────────────────────────────
out = r"c:\Users\Usuario\Pictures\PYTHON\Backend-Chatbot\docs\gerencia\CONFIA_Presentacion_Gerencia.pptx"
prs.save(out)
print(f"Guardado: {out}")
